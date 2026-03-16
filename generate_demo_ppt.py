"""
Generate demo PowerPoint for UPI AI Hackathon project.
Requires: pip install python-pptx
Run: python generate_demo_ppt.py
Output: UPI_AI_Hackathon_Demo.pptx (in project root)
"""
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Please install python-pptx: pip install python-pptx")
    raise

# Output path
ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "UPI_AI_Hackathon_Demo.pptx"

# Theme colors (NPCI / UPI style)
BLUE = RGBColor(0x1E, 0x3A, 0x5F)   # Dark blue
TEAL = RGBColor(0x0D, 0x94, 0x88)   # Teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)


def set_title(shape, text, size_pt=28, bold=True):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = BLUE


def set_body(shape, text, size_pt=14):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for line in text.strip().split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size_pt)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)


def add_bullet_slide(prs, title_text, bullets, subbullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(0.5)
    top = Inches(0.6)
    w = Inches(9)
    h = Inches(0.8)
    title = slide.shapes.add_textbox(left, top, w, h)
    set_title(title, title_text, 24)
    body_top = Inches(1.5)
    body = slide.shapes.add_textbox(left, body_top, w, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        if i:
            p.text = ""
        p.text = (p.text or "") + "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    return slide


def add_content_slide(prs, title_text, body_text, font_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.6)
    w = Inches(9)
    title = slide.shapes.add_textbox(left, top, w, Inches(0.8))
    set_title(title, title_text, 24)
    body = slide.shapes.add_textbox(left, Inches(1.5), w, Inches(5.5))
    set_body(body, body_text, font_size)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    box = s1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "UPI AI Hackathon – Project Demo"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    sub = s1.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.6))
    sub.text_frame.paragraphs[0].text = "Phase 1: UPI Switch Simulation  |  Phase 2: AI Agents for Spec Change"
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = TEAL
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 2: Problem Statement (Phase 1 + Phase 2) ---
    add_bullet_slide(prs, "Problem Statement – NPCI AI Hackathon", [
        "Phase 1: Build end-to-end UPI switch simulation with Payer PSP, Payee PSP, Remitter Bank, Beneficiary Bank, NPCI UPI Switch.",
        "Phase 2: AI agents for all five parties to automate change management: NPCI issues spec changes → Bank/PSP agents interpret, update code/tests/docs, signal deploy readiness.",
        "Key rule: All API calls asynchronous (ACK + TxnId; final status via callback/poll).",
        "Compliance: XSD validation, XML parsing at each hop, transaction logs, Kafka-like event stream, Redis/key DB.",
    ])

    # --- Slide 3: Phase 1 – Component Responsibilities ---
    add_content_slide(prs, "Phase 1 – Component Responsibilities", """
• API Client / UPI App (dummy)  →  Triggers XML requests (Postman/CLI/UI)
• Payee PSP (dummy)  →  Forwards XML to switch; address/auth resolution
• UPI Switch (core)  →  Parses/validates XML, routes, logs, TxnId/RRN, coordinates flow
• Payer PSP (dummy)  →  Validates VPA, PIN (CredBlock), device checks (dummy)
• Remitter Bank (dummy)  →  Dummy CBS: validate balance, debit, return status
• Beneficiary Bank (dummy)  →  Dummy CBS: credit account, return status

Switch does not store balances; only banks do. Switch routes, validates, logs, returns statuses.
    """, 13)

    # --- Slide 4: How UPI Pay Transaction Works (Flow) ---
    add_content_slide(prs, "Phase 1 – UPI Pay Transaction Flow", """
1. User enters payee VPA / scans QR in UPI app
2. Address validation → Payee PSP via switch
3. Payee PSP returns resolved details to Payer PSP via switch
4. Payer takes user auth (UPI PIN), confirms in app
5. Payer PSP sends Pay request to UPI switch
6. Switch → Payee PSP (auth) to resolve beneficiary account
7. Switch → Remitter Bank (debit request)
8. Remitter Bank validates PIN (cred block), debits account, returns status
9. Switch → Beneficiary Bank (credit request)
10. Beneficiary Bank credits account, returns status
11. Final status flows back to both PSPs and the app
    """, 12)

    # --- Slide 5: Phase 2 – What to Build ---
    add_bullet_slide(prs, "Phase 2 – AI Agents for Spec Change", [
        "NPCI Switch (Core AI Agent): Creates change Manifest, dispatches to all parties.",
        "Remitter Bank / Beneficiary Bank (Core AI Agents): Update debit/credit CBS logic and validators per change.",
        "Coordinator/Orchestrator: Tracks change status across parties (RECEIVED → APPLIED → TESTED → READY).",
        "PSPs and app can be hard-coded from Phase 1; ready for change after NPCI initiation.",
        "Agents communicate using A2A protocol; signed manifests (bonus).",
    ])

    # --- Slide 6: High-Level Architecture (Our Implementation) ---
    add_content_slide(prs, "Our Implementation – High-Level Architecture", """
Phase 1 (javacoderepo):
  • UPISim          – NPCI UPI Switch (Java, port 8080); routing, ReqPay, logs, DB
  • PayerPSP        – Payer PSP (Java)
  • PayeePSP        – Payee PSP (Java)
  • RemitterBank    – Remitter Bank with dummy CBS
  • BeneficiaryBank – Beneficiary Bank with dummy CBS

Phase 2 (Allagents):
  • aicode          – NPCI AI Agent (FastAPI): change requests, manifest creation, signing, broadcast, XSD/OpenAPI, orchestrator
  • Payer-agent     – Payer PSP AI agent (port 9001): receive manifest, patch code, tests, report status
  • Payee-agent     – Payee PSP AI agent (port 9004)
  • Remitter-agent  – Remitter Bank AI agent (port 9003)
  • Beneficiary-agent – Beneficiary Bank AI agent (port 9002)
    """, 12)

    # --- Slide 7: Phase 2 – Change Propagation Flow ---
    add_content_slide(prs, "Phase 2 – Change Propagation Flow", """
1. Product/Developer submits change request (e.g. new purpose code, new API) in Admin Portal
2. NPCI agent generates spec/XSD, creates manifest (optionally signed with RSA-SHA256)
3. Manifest broadcast to all partners (POST to each agent’s /agent/manifest)
4. Each agent: RECEIVED → applies patch (code/tests/XSD) → TESTS_READY → TESTED → APPROVE → READY
5. Agents notify orchestrator via POST /orchestrator/a2a/status { changeId, agent, status }
6. Status board: RECEIVED → APPLIED → TESTED → READY (per change, per party)
7. When all READY → change can be rolled out; transactions use new specs
    """, 12)

    # --- Slide 8: Orchestrator & A2A ---
    add_content_slide(prs, "Orchestrator & A2A Protocol", """
• Orchestrator state: changeId → { Payer, Payee, Remitter, Beneficiary } → status
• Status values: RECEIVED | APPLIED | TESTED | TESTS_READY | READY (no downgrade)
• When NPCI broadcasts manifest, delivery is recorded as RECEIVED for that party
• Parties POST status updates: /orchestrator/a2a/status
• Frontend (Admin Portal) shows Status Center / Partners tab with per-change, per-party status
• Signed manifest: NPCI signs with private key; partners verify with public key (GET /npciswitch/keys/public)
    """, 13)

    # --- Slide 9: Admin Portal (Demo UI) ---
    add_content_slide(prs, "Admin Portal – Demo UI", """
• Roles: Product, Developer, Admin (different tabs and themes)
• Product: Submit change requests (API, change type, field additions, XSD types)
• Developer: Review change requests, approve, trigger manifest creation, broadcast/send to partner
• Access Management: Partners, manifest delivery status, download signed manifest
• Status Center: View orchestrator status (RECEIVED → APPLIED → TESTED → READY) per change
• Each bank/PSP agent has its own UI: receive manifest, run process, approve diff, mark TESTED/READY
    """, 13)

    # --- Slide 10: Minimum Deliverables Checklist ---
    add_bullet_slide(prs, "Minimum Deliverables – Addressed", [
        "Phase 1: Working demo of UPI Pay transaction (App → PSPs → Switch → Banks); all components as endpoints.",
        "Async API messaging: ACK + callback/poll; XML parsing + XSD validation; transaction log (TxnId, timestamps); final status in XML.",
        "Phase 2: Two+ running agents with message API; end-to-end demo (NPCI manifest → agents update code/tests/docs).",
        "Saved code diffs, unit test outputs, XSD updates; status board RECEIVED → APPLIED → TESTED → READY.",
        "NPCI and Bank agents running; A2A protocol for status updates. Bonus: signed manifests, transaction visibility in dashboard.",
    ])

    # --- Slide 11: How to Run Demo ---
    add_content_slide(prs, "How to Run the Demo", """
Phase 1 – UPI transaction:
  • Start: PayerPSP, PayeePSP, RemitterBank, BeneficiaryBank, UPISim (see javacoderepo/README.md)
  • DB: Create upisim DB, psp_bank and transaction_logs tables
  • Trigger Pay: GET http://localhost:8080/upi/Pay?payerVpa=toshu@pyr&payeeVpa=kalai@pye&amount=900.00

Phase 2 – Spec change:
  • Start aicode (e.g. port 8000), Payer-agent (9001), Payee-agent (9004), Remitter-agent (9003), Beneficiary-agent (9002)
  • Open Admin Portal → Product: submit change → Developer: approve → Create manifest → Broadcast
  • Open each agent UI to process manifest and report status; view Status Center on NPCI
    """, 11)

    # --- Slide 12: Thank You ---
    s_end = prs.slides.add_slide(prs.slide_layouts[6])
    box = s_end.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    q = s_end.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.6))
    q.text_frame.paragraphs[0].text = "Q & A"
    q.text_frame.paragraphs[0].font.size = Pt(22)
    q.text_frame.paragraphs[0].font.color.rgb = TEAL
    q.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
